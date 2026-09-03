from pathlib import Path
from pptx import Presentation

deck = Presentation(Path(__file__).resolve().parents[1] / "SIH26145_SentinelAI_OS_Idea_Presentation.pptx")
assert len(deck.slides) == 6, f"Expected 6 slides, got {len(deck.slides)}"
required = ("SIH26145", "Unidirectional Defense", "Technical Approach", "Feasibility", "Impact", "Traceability")
all_text = "\n".join(shape.text for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text"))
missing = [item for item in required if item.lower() not in all_text.lower()]
assert not missing, f"Missing: {missing}"
print(f"Verified {len(deck.slides)} slides; required SIH content is present.")
