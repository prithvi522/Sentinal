from pathlib import Path
from copy import deepcopy

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = Path(r"C:\Users\prith\Downloads\SIH2026-IDEA-Presentation-Format.pptx")
OUTPUT = ROOT / "SIH26145_SentinelAI_OS_Idea_Presentation.pptx"

NAVY = RGBColor(5, 13, 28)
PANEL = RGBColor(10, 27, 48)
CYAN = RGBColor(30, 210, 235)
LIME = RGBColor(164, 230, 53)
WHITE = RGBColor(236, 248, 255)
MUTED = RGBColor(169, 192, 211)
AMBER = RGBColor(251, 191, 36)
ROSE = RGBColor(251, 113, 133)


def remove_slide(prs, index):
    slide_id = prs.slides._sldIdLst[index]
    rel_id = slide_id.rId
    prs.part.drop_rel(rel_id)
    del prs.slides._sldIdLst[index]


def clear_text(slide):
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape.text_frame.clear()


def rect(slide, x, y, w, h, color=PANEL, transparency=0, line=None, radius=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency
    shape.line.color.rgb = line if line else color
    if radius:
        shape.adjustments[0] = 0.08
    return shape


def text(slide, value, x, y, w, h, size=16, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="Aptos", margin=0.06):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame; frame.clear(); frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]; paragraph.alignment = align
    run = paragraph.add_run(); run.text = value
    run.font.name = font; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
    return box


def bullets(slide, items, x, y, w, h, size=13.5, color=WHITE, accent=CYAN):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame; frame.clear(); frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.12)
    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = item; p.level = 0; p.font.name = "Aptos"; p.font.size = Pt(size); p.font.color.rgb = color
        p.space_after = Pt(8); p.bullet = True
    return box


def header(slide, number, title, eyebrow):
    text(slide, f"{number:02d}", 0.55, 0.34, 0.5, 0.3, 12, CYAN, True)
    text(slide, eyebrow.upper(), 1.08, 0.32, 4.7, 0.3, 10, MUTED, True)
    text(slide, title, 0.55, 0.67, 11.95, 0.55, 25, WHITE, True)
    rect(slide, 0.55, 1.28, 12.1, 0.02, CYAN)


def footer(slide):
    text(slide, "SENTINELAI OS  •  SIH26145  •  PASSIVE DEFENCE ONLY", 0.55, 7.12, 8.0, 0.2, 8.5, MUTED, True)
    text(slide, "Team: [REGISTERED TEAM NAME]", 9.3, 7.12, 3.35, 0.2, 8.5, MUTED, True, PP_ALIGN.RIGHT)


def pipeline(slide, labels, x, y, width):
    unit = width / len(labels)
    for index, label in enumerate(labels):
        left = x + index * unit
        rect(slide, left, y, unit - 0.1, 0.58, PANEL, line=CYAN, radius=True)
        text(slide, label, left + 0.06, y + 0.14, unit - 0.22, 0.25, 10, WHITE, True, PP_ALIGN.CENTER)
        if index < len(labels) - 1:
            text(slide, "›", left + unit - 0.08, y + 0.11, 0.18, 0.28, 19, LIME, True, PP_ALIGN.CENTER)


def build():
    prs = Presentation(TEMPLATE)
    while len(prs.slides) > 6:
        remove_slide(prs, len(prs.slides) - 1)
    for slide in prs.slides:
        clear_text(slide)

    # 1 — title
    s = prs.slides[0]
    rect(s, 0, 0, 13.333, 7.5, NAVY)
    rect(s, 0, 0, 0.18, 7.5, CYAN)
    text(s, "SMART INDIA HACKATHON 2026", 0.72, 0.5, 6.8, 0.35, 14, CYAN, True)
    text(s, "SENTINELAI OS", 0.72, 1.05, 7.6, 0.65, 37, WHITE, True)
    text(s, "AI-Powered Passive Traffic Intelligence", 0.72, 1.74, 8.2, 0.38, 21, LIME, True)
    text(s, "Unidirectional Defense Engine for Critical Infrastructure", 0.72, 2.2, 8.4, 0.35, 16, MUTED)
    pipeline(s, ["ONE-WAY INGEST", "ANALYZE", "SCORE", "EXPLAIN", "ALERT"], 0.72, 3.08, 11.8)
    rect(s, 0.72, 4.15, 11.8, 1.48, PANEL, line=CYAN, radius=True)
    text(s, "SIH26145  •  AI-Based Detection of Cyber Threats in Unidirectional IP Traffic", 1.02, 4.47, 11.15, 0.3, 18, WHITE, True, PP_ALIGN.CENTER)
    text(s, "Read-only ingest  |  No return path  |  Metadata-only TLS/QUIC analysis  |  Streaming SOC alerts", 1.02, 4.93, 11.15, 0.3, 13, MUTED, False, PP_ALIGN.CENTER)
    text(s, "Problem Statement ID: SIH26145\nTheme: Cyber Security\nPS Category: Software", 0.72, 6.12, 4.1, 0.7, 12, WHITE)
    text(s, "Team ID: [TEAM ID]\nTeam Name: [REGISTERED TEAM NAME]\nInstitute: [INSTITUTE NAME]", 8.3, 6.12, 4.2, 0.7, 12, WHITE, False, PP_ALIGN.RIGHT)

    # 2 — solution
    s = prs.slides[1]; rect(s, 0, 0, 13.333, 7.5, NAVY); header(s, 2, "Proposed Solution", "Flagship SIH26145 capability")
    rect(s, .55, 1.6, 5.8, 4.95, PANEL, line=CYAN, radius=True)
    text(s, "THE PROBLEM", .82, 1.87, 2.3, .28, 13, AMBER, True)
    bullets(s, ["Critical networks use passive mirroring or data diodes; the monitoring enclave cannot query, handshake with, or control production.", "Traditional tools often depend on active scanning, payload access, or a return path—unsafe and out of scope.", "SOC teams need near-real-time, explainable intelligence from one-way traffic only."], .78, 2.25, 5.25, 3.7)
    rect(s, 6.62, 1.6, 6.05, 4.95, PANEL, line=LIME, radius=True)
    text(s, "OUR SOLUTION", 6.9, 1.87, 2.7, .28, 13, LIME, True)
    bullets(s, ["SentinelAI OS transforms observed PCAP, passive TAP, and normalized flow records into evidence-backed alerts.", "Hybrid local detection: configurable rules + behavioral baseline + optional local anomaly ML; core detection remains offline.", "Detects DDoS, C2 beaconing, DGA/DNS tunnelling, encrypted-session anomalies, reconnaissance, and data exfiltration.", "Outputs confidence, severity, features, and recommended human investigation—never automatic mitigation."], 6.84, 2.25, 5.48, 3.8)
    text(s, "UNIQUE VALUE: monitoring that cannot become a pivot into the protected network.", .8, 6.77, 11.7, .27, 12, CYAN, True, PP_ALIGN.CENTER); footer(s)

    # 3 — tech
    s = prs.slides[2]; rect(s, 0, 0, 13.333, 7.5, NAVY); header(s, 3, "Technical Approach", "Passive streaming architecture")
    pipeline(s, ["PCAP / TAP / FLOW", "NORMALIZE", "FEATURES", "RULES + ML", "EVIDENCE", "SOC"], .58, 1.62, 12.1)
    rect(s, .58, 2.5, 4.0, 3.85, PANEL, line=CYAN, radius=True); text(s, "PASSIVE INGEST", .84, 2.78, 2.5, .25, 13, CYAN, True)
    bullets(s, ["Scapy metadata extraction: IPs, ports, protocols, timestamps, TCP flags, DNS queries.", "PCAP/replay/live receive-only ingest and normalized NetFlow/IPFIX/sFlow-style flow API.", "Bounded queues, configurable workers, backpressure and WebSocket updates."], .76, 3.13, 3.55, 2.7, 12.2)
    rect(s, 4.72, 2.5, 4.0, 3.85, PANEL, line=LIME, radius=True); text(s, "DETECTION + ML", 4.98, 2.78, 2.8, .25, 13, LIME, True)
    bullets(s, ["Rates, entropy, DNS n-gram signals, timing periodicity, TLS metadata, fan-out and flow asymmetry.", "DDoS • C2 • DGA • DNS tunnel • TLS anomaly • recon • exfiltration.", "Optional local Isolation Forest; explicit heuristic/statistical fallback if no trained model exists."], 4.9, 3.13, 3.55, 2.7, 12.2)
    rect(s, 8.86, 2.5, 3.82, 3.85, PANEL, line=AMBER, radius=True); text(s, "TECH STACK", 9.12, 2.78, 2.3, .25, 13, AMBER, True)
    bullets(s, ["React + Tailwind + WebSockets", "FastAPI + Pydantic + SQLAlchemy", "PostgreSQL / SQLite dev compatibility", "Docker Compose; configurable data-diode NIC allow-list"], 9.04, 3.13, 3.35, 2.7, 12.2); footer(s)

    # 4 feasibility
    s = prs.slides[3]; rect(s, 0, 0, 13.333, 7.5, NAVY); header(s, 4, "Feasibility and Viability", "Safe deployment roadmap")
    rect(s, .58, 1.58, 5.8, 4.95, PANEL, line=LIME, radius=True); text(s, "WHY IT IS FEASIBLE", .86, 1.86, 3.3, .28, 13, LIME, True)
    bullets(s, ["Built as an extension of the existing SentinelAI OS: reuse authentication, FastAPI APIs, SQLAlchemy, dashboard, WebSockets and Docker.", "Works from PCAP/replay before a hardware data diode is available; supports receive-only live capture later.", "No external AI, cloud API, or threat-intelligence service is required for core operation.", "Dedicated flow, alert, baseline and benchmark persistence already supports incremental validation."], .78, 2.22, 5.28, 3.78, 12.5)
    rect(s, 6.62, 1.58, 6.05, 4.95, PANEL, line=CYAN, radius=True); text(s, "RISK → CONTROL", 6.9, 1.86, 2.5, .28, 13, CYAN, True)
    bullets(s, ["Traffic burst → bounded queue, worker scaling, queue-depth telemetry and measured benchmark runner.", "False positives → multi-signal fusion, rolling baselines and analyst-visible evidence.", "Sensitive data → metadata-first processing; no payload persistence or TLS/QUIC decryption.", "Monitoring pivot risk → no transmit, probe, block, firewall, quarantine or command path is implemented.", "ML unavailable → rules/statistics continue; synthetic data is labelled demo-only, not accuracy evidence."], 6.84, 2.22, 5.48, 3.9, 12.1); footer(s)

    # 5 impact
    s = prs.slides[4]; rect(s, 0, 0, 13.333, 7.5, NAVY); header(s, 5, "Impact and Benefits", "Critical-infrastructure SOC outcomes")
    cards = [("SAFE BY DESIGN", "A one-way monitoring architecture removes the analytics platform as a route back into production.", CYAN), ("ACTIONABLE INTELLIGENCE", "Evidence-rich alerts shorten triage: source, destination, threat, confidence, severity and observed features.", LIME), ("PRIVACY PRESERVING", "Encrypted-session detection uses timing, size and fingerprint metadata—payload decryption remains disabled.", AMBER), ("DEPLOYABLE IN STAGES", "Replay PCAP → validate with analysts → passive TAP/data-diode feed → normalized flow-collector integration.", ROSE)]
    for index, (title, body, color) in enumerate(cards):
        x = .58 + (index % 2) * 6.13; y = 1.65 + (index // 2) * 2.25
        rect(s, x, y, 5.83, 1.82, PANEL, line=color, radius=True)
        text(s, title, x + .24, y + .22, 5.2, .24, 13, color, True)
        text(s, body, x + .24, y + .62, 5.22, .82, 12.3, WHITE)
    text(s, "Target sectors: Power • Oil & Gas • Transport • Telecom • Defence • Government • Manufacturing", .75, 6.45, 11.85, .3, 13, CYAN, True, PP_ALIGN.CENTER); footer(s)

    # 6 research / traceability
    s = prs.slides[5]; rect(s, 0, 0, 13.333, 7.5, NAVY); header(s, 6, "Research, Validation and Traceability", "Standards-aligned prototype")
    rect(s, .58, 1.55, 6.02, 4.95, PANEL, line=CYAN, radius=True); text(s, "SIH REQUIREMENT → IMPLEMENTATION", .84, 1.84, 4.2, .27, 13, CYAN, True)
    bullets(s, ["Read-only / no return path → data-diode-safe passive ingest; no control or mitigation subsystem.", "Streaming → bounded queues, flow workers, WebSocket flow/alert/metric events.", "No payload decryption → TLS/QUIC metadata-only analysis.", "Threat coverage → DDoS, C2, DGA/DNS tunnel, encrypted anomaly, recon, exfiltration detectors.", "Standard alert → Pydantic schema + PostgreSQL persistence + explainable evidence.", "Throughput → benchmark endpoint records actual local flows/sec; no invented performance claim."], .78, 2.2, 5.55, 3.85, 11.7)
    rect(s, 6.82, 1.55, 5.86, 4.95, PANEL, line=LIME, radius=True); text(s, "RESEARCH + VALIDATION PLAN", 7.08, 1.84, 3.5, .27, 13, LIME, True)
    bullets(s, ["NIST SP 800-94 — Guide to Intrusion Detection and Prevention Systems.", "CISA ICS defense-in-depth guidance — unidirectional gateways/data diodes.", "MITRE ATT&CK: T1498 DoS, T1046 Network Service Scanning, T1071.004 DNS.", "Validation: safe synthetic scenarios + curated PCAP replay + analyst review + measured benchmark.", "Future work: train and validate supervised models on approved labelled datasets; tune thresholds per protected environment."], 7.02, 2.2, 5.38, 3.85, 11.7)
    text(s, "Final security pledge: OBSERVE → EXTRACT → ANALYZE → DETECT → SCORE → EXPLAIN → ALERT.  NEVER probe, connect back, inject, decrypt, block or control.", .78, 6.68, 11.8, .3, 11.2, WHITE, True, PP_ALIGN.CENTER); footer(s)

    prs.save(OUTPUT)
    print(OUTPUT)

if __name__ == "__main__":
    build()
